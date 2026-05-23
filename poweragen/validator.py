from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from .utils import iter_shapes, now_iso, shape_text


UNRESOLVED_TOKENS = ["{{", "}}", "<<", ">>", "TODO", "PLACEHOLDER", "Lorem ipsum"]


class PptxValidator:
    def validate(
        self,
        pptx_path: Path,
        slide_specs: dict[str, Any],
        layout_plan: dict[str, Any],
    ) -> dict[str, Any]:
        issues: list[dict[str, Any]] = []
        try:
            prs = Presentation(str(pptx_path))
        except Exception as exc:
            return {
                "schema_version": "0.1",
                "generated_at": now_iso(),
                "status": "fail",
                "issues": [
                    {
                        "slide_no": 0,
                        "severity": "error",
                        "code": "PPTX_OPEN_FAILED",
                        "message": str(exc),
                        "shape_id": None,
                    }
                ],
                "metrics": {
                    "slide_count": 0,
                    "editable_text_ratio": 0.0,
                    "template_similarity_estimate": 0.0,
                },
            }

        editable_text_shapes = 0
        total_shapes = 0
        for slide_no, slide in enumerate(prs.slides, start=1):
            visible_text = []
            slide_shapes = list(iter_shapes(slide.shapes))
            total_shapes += len(slide_shapes)
            for shape in slide_shapes:
                text = shape_text(shape)
                if getattr(shape, "has_text_frame", False):
                    editable_text_shapes += 1
                if text:
                    visible_text.append(text)
                    for token in UNRESOLVED_TOKENS:
                        if token in text:
                            issues.append(
                                {
                                    "slide_no": slide_no,
                                    "severity": "error",
                                    "code": "UNRESOLVED_PLACEHOLDER",
                                    "message": f"Unresolved token {token} found.",
                                    "shape_id": str(shape.shape_id),
                                }
                            )
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt < 9:
                                issues.append(
                                    {
                                        "slide_no": slide_no,
                                        "severity": "warn",
                                        "code": "FONT_TOO_SMALL",
                                        "message": "A text run is smaller than 9pt.",
                                        "shape_id": str(shape.shape_id),
                                    }
                                )
            if not visible_text:
                issues.append(
                    {
                        "slide_no": slide_no,
                        "severity": "error",
                        "code": "EMPTY_SLIDE_TEXT",
                        "message": "Slide has no editable visible text.",
                        "shape_id": None,
                    }
                )

        expected_count = len(slide_specs["specs"])
        if len(prs.slides) != expected_count:
            issues.append(
                {
                    "slide_no": 0,
                    "severity": "error",
                    "code": "SLIDE_COUNT_MISMATCH",
                    "message": f"Expected {expected_count} slides but found {len(prs.slides)}.",
                    "shape_id": None,
                }
            )

        strategy_counts: dict[str, int] = {}
        for item in layout_plan["plan"]:
            strategy = item["rendering_strategy"]
            strategy_counts[strategy] = strategy_counts.get(strategy, 0) + 1
        clone_ratio = strategy_counts.get("clone_fill", 0) / max(1, expected_count)
        error_count = sum(1 for issue in issues if issue["severity"] == "error")
        warn_count = sum(1 for issue in issues if issue["severity"] == "warn")
        status = "fail" if error_count else ("warn" if warn_count else "pass")

        return {
            "schema_version": "0.1",
            "generated_at": now_iso(),
            "status": status,
            "issues": issues,
            "metrics": {
                "slide_count": len(prs.slides),
                "editable_text_ratio": round(editable_text_shapes / max(1, total_shapes), 3),
                "template_similarity_estimate": round(min(1.0, 0.55 + clone_ratio * 0.45), 3),
                "strategy_counts": strategy_counts,
            },
        }
