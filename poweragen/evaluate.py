from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation

from .agent import PowerAgen
from .utils import write_json


DEFAULT_EVAL_PROMPT = (
    "Create an 8-slide graduate-level report about AI-assisted course planning. "
    "Cover background, core problems, findings, action path, execution rhythm, and next steps."
)


def run_evaluation(
    template_path: Path,
    run_dir: Path,
    prompt: str = DEFAULT_EVAL_PROMPT,
    slide_count: int = 8,
    require_strategy_b: bool = True,
) -> dict[str, Any]:
    result = PowerAgen().run(
        template_path=template_path,
        prompt=prompt,
        slide_count=slide_count,
        run_dir=run_dir,
        visual_strictness="adaptive",
    )
    checks: list[dict[str, Any]] = []

    def check(name: str, passed: bool, detail: str) -> None:
        checks.append({"name": name, "passed": bool(passed), "detail": detail})

    expected_artifacts = [
        "user_intent.json",
        "template_schema.json",
        "deck_outline.json",
        "slide_specs.json",
        "validation_report.json",
    ]
    for artifact in expected_artifacts:
        check(
            f"artifact:{artifact}",
            (run_dir / artifact).exists(),
            f"{artifact} exists under run directory",
        )

    pptx = Presentation(result["output_path"])
    check("pptx_openable", True, "output.pptx opens with python-pptx")
    check("slide_count", len(pptx.slides) == slide_count, f"found {len(pptx.slides)} slides")

    layout_plan = json.loads((run_dir / "slide_layout_plan.json").read_text(encoding="utf-8"))
    validation = json.loads((run_dir / "validation_report.json").read_text(encoding="utf-8"))
    strategies = [item["rendering_strategy"] for item in layout_plan["plan"]]
    clone_ratio = strategies.count("clone_fill") / max(1, len(strategies))
    check("strategy_a_ratio", clone_ratio >= 0.7, f"clone_fill ratio is {clone_ratio:.2f}")
    if require_strategy_b:
        check("strategy_b_present", "template_extend" in strategies, "at least one slide uses Strategy B")
    else:
        check("strategy_b_optional", True, "Strategy B is optional for this template family")
    check("validation_status", validation["status"] in {"pass", "warn"}, f"status is {validation['status']}")
    check(
        "no_unresolved_placeholders",
        not any(issue["code"] == "UNRESOLVED_PLACEHOLDER" for issue in validation["issues"]),
        "validator found no unresolved placeholder tokens",
    )

    passed = all(item["passed"] for item in checks)
    report = {
        "schema_version": "0.1",
        "status": "pass" if passed else "fail",
        "run_result": result,
        "checks": checks,
    }
    write_json(run_dir / "evaluation_report.json", report)
    return report


def run_evaluation_suite(
    templates: list[Path],
    run_dir: Path,
    prompt: str = DEFAULT_EVAL_PROMPT,
    slide_count: int = 8,
) -> dict[str, Any]:
    run_dir.mkdir(parents=True, exist_ok=True)
    deck_reports = []
    strategy_b_seen = False
    for template in templates:
        child_dir = run_dir / template.stem
        report = run_evaluation(
            template_path=template,
            run_dir=child_dir,
            prompt=prompt,
            slide_count=slide_count,
            require_strategy_b=False,
        )
        deck_reports.append({"template": str(template), "report": report})
        strategy_b_seen = strategy_b_seen or _report_has_strategy_b(child_dir)

    checks = [
        {
            "name": "all_decks_pass",
            "passed": all(item["report"]["status"] == "pass" for item in deck_reports),
            "detail": f"{len(deck_reports)} template families evaluated",
        },
        {
            "name": "suite_strategy_b_present",
            "passed": strategy_b_seen,
            "detail": "at least one evaluated deck uses Strategy B",
        },
    ]
    passed = all(item["passed"] for item in checks)
    report = {
        "schema_version": "0.1",
        "status": "pass" if passed else "fail",
        "checks": checks,
        "deck_reports": deck_reports,
    }
    write_json(run_dir / "evaluation_suite_report.json", report)
    return report


def _report_has_strategy_b(run_dir: Path) -> bool:
    plan_path = run_dir / "slide_layout_plan.json"
    if not plan_path.exists():
        return False
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    return any(item["rendering_strategy"] == "template_extend" for item in plan.get("plan", []))


if __name__ == "__main__":
    report = run_evaluation(Path("test/courseplan_test.pptx"), Path("runs/eval-smoke"))
    raise SystemExit(0 if report["status"] == "pass" else 1)
