from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from .utils import EMU_PER_INCH, iter_shapes


class PptxRenderer:
    def render(self, template_path: Path, slide_specs: dict[str, Any], output_path: Path) -> Path:
        prs = Presentation(str(template_path))
        original_count = len(prs.slides)
        rendered_slides = []

        for spec in slide_specs["specs"]:
            template_ref = spec["template_ref"]
            if not template_ref:
                raise ValueError("MVP renderer requires a template_ref for every slide.")
            slide = self._clone_slide(prs, prs.slides[template_ref - 1])
            if spec["rendering_strategy"] == "template_extend":
                self._render_grid_extend(prs, slide, spec)
            else:
                self._render_clone_fill(slide, spec)
            rendered_slides.append(slide)

        for _ in range(original_count):
            self._delete_slide(prs, 0)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        return output_path

    def _clone_slide(self, prs: Presentation, source_slide: Any) -> Any:
        dest = prs.slides.add_slide(prs.slide_layouts[6])
        relmap: dict[str, str] = {}
        for rel in source_slide.part.rels.values():
            if "notesSlide" in rel.reltype or "slideLayout" in rel.reltype:
                continue
            new_rid = dest.part.rels._add_relationship(rel.reltype, rel._target, rel.is_external)
            relmap[rel.rId] = new_rid

        for shape in source_slide.shapes:
            new_el = deepcopy(shape.element)
            for el in new_el.iter():
                for attr, value in list(el.attrib.items()):
                    if value in relmap:
                        el.set(attr, relmap[value])
            dest.shapes._spTree.insert_element_before(new_el, "p:extLst")
        return dest

    def _delete_slide(self, prs: Presentation, index: int) -> None:
        slide_id = prs.slides._sldIdLst[index]
        prs.part.drop_rel(slide_id.rId)
        prs.slides._sldIdLst.remove(slide_id)

    def _render_clone_fill(self, slide: Any, spec: dict[str, Any]) -> None:
        bindings = (spec.get("recipe_ref") or {}).get("bindings", {})
        values = {item["slot_id"]: item["value"] for item in spec.get("slot_values", [])}
        filled_shape_ids: set[int] = set()

        for slot_id, value in values.items():
            binding = bindings.get(slot_id)
            if not binding:
                continue
            shape = self._find_shape(slide, int(binding["shape_id"]))
            if shape is None:
                continue
            self._set_text(shape, value, is_title=slot_id == "title")
            filled_shape_ids.add(int(binding["shape_id"]))

        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and int(shape.shape_id) not in filled_shape_ids:
                self._set_text(shape, " ")

    def _render_grid_extend(self, prs: Presentation, slide: Any, spec: dict[str, Any]) -> None:
        self._render_clone_fill(slide, spec)
        tokens = spec.get("design_tokens") or {}
        spacing = tokens.get("spacing_emu", {})
        colors = tokens.get("colors", {})
        content_top = int(spacing.get("content_top") or prs.slide_height * 0.3)

        for shape in list(slide.shapes):
            if int(shape.top) >= content_top:
                self._remove_shape(shape)

        items = []
        for group in spec.get("item_groups", []):
            if group.get("slot_id") == "grid_items":
                items = group.get("items", [])
                break
        target_count = (spec.get("extension_params") or {}).get("target_count") or len(items)
        items = items[:target_count]
        if not items:
            return

        self._draw_grid(slide, prs.slide_width, prs.slide_height, items, colors, spacing)

    def _draw_grid(
        self,
        slide: Any,
        slide_width: int,
        slide_height: int,
        items: list[dict[str, str]],
        colors: dict[str, str],
        spacing: dict[str, int],
    ) -> None:
        count = len(items)
        cols = 3 if count <= 6 else 4
        if count <= 4:
            cols = count
        rows = (count + cols - 1) // cols
        margin_x = int(spacing.get("margin_left") or slide_width * 0.07)
        top = int(spacing.get("content_top") or slide_height * 0.3)
        gap = int(spacing.get("grid_gap") or slide_width * 0.018)
        available_w = slide_width - margin_x * 2
        available_h = slide_height - top - int(0.55 * EMU_PER_INCH)
        card_w = int((available_w - gap * (cols - 1)) / cols)
        card_h = int((available_h - gap * (rows - 1)) / rows)
        primary = self._rgb(colors.get("primary", "#174A7C"))
        accent = self._rgb(colors.get("accent", "#F2B705"))

        for idx, item in enumerate(items):
            row, col = divmod(idx, cols)
            left = margin_x + col * (card_w + gap)
            y = top + row * (card_h + gap)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)
            card.line.color.rgb = primary
            card.line.width = Pt(1.1)

            badge = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left + int(0.16 * EMU_PER_INCH),
                y + int(0.14 * EMU_PER_INCH),
                int(0.45 * EMU_PER_INCH),
                int(0.45 * EMU_PER_INCH),
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = accent
            badge.line.color.rgb = accent
            self._set_text(badge, item.get("label", f"{idx + 1:02d}"), is_title=True, size_pt=12)

            title_box = slide.shapes.add_textbox(
                left + int(0.72 * EMU_PER_INCH),
                y + int(0.11 * EMU_PER_INCH),
                card_w - int(0.88 * EMU_PER_INCH),
                int(0.5 * EMU_PER_INCH),
            )
            self._set_text(title_box, item.get("title", ""), is_title=True, size_pt=15)

            body_box = slide.shapes.add_textbox(
                left + int(0.22 * EMU_PER_INCH),
                y + int(0.78 * EMU_PER_INCH),
                card_w - int(0.44 * EMU_PER_INCH),
                card_h - int(0.98 * EMU_PER_INCH),
            )
            self._set_text(body_box, item.get("body", ""), size_pt=12)

    def _find_shape(self, slide: Any, shape_id: int) -> Any | None:
        for shape in iter_shapes(slide.shapes):
            if int(shape.shape_id) == shape_id:
                return shape
        return None

    def _set_text(
        self,
        shape: Any,
        text: str,
        is_title: bool = False,
        size_pt: int | None = None,
    ) -> None:
        if not getattr(shape, "has_text_frame", False):
            return
        text = (text or " ").strip() or " "
        shape.text_frame.clear()
        lines = text.splitlines() or [" "]
        first = shape.text_frame.paragraphs[0]
        first.text = lines[0]
        for line in lines[1:]:
            para = shape.text_frame.add_paragraph()
            para.text = line
        final_size = size_pt or (18 if is_title else 12)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(max(14 if is_title else 12, final_size))

    def _remove_shape(self, shape: Any) -> None:
        element = shape.element
        parent = element.getparent()
        if parent is not None:
            parent.remove(element)

    def _rgb(self, value: str) -> RGBColor:
        clean = value.strip().lstrip("#")
        if len(clean) != 6:
            return RGBColor(31, 41, 51)
        return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))
