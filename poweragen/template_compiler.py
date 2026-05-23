from __future__ import annotations

import re
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from .utils import EMU_PER_INCH, iter_shapes, now_iso, shape_text


COMMON_LAYOUTS = [
    "cover",
    "divider",
    "agenda",
    "content",
    "comparison",
    "grid",
    "sequence",
    "chart",
    "quote",
    "closing",
]


class TemplateCompiler:
    def compile(self, template_path: Path) -> dict[str, Any]:
        prs = Presentation(str(template_path))
        slide_recipes: list[dict[str, Any]] = []
        catalog: list[dict[str, Any]] = []
        all_colors: list[str] = []
        all_fonts: list[str] = []
        all_sizes: list[int] = []

        for index, slide in enumerate(prs.slides, start=1):
            text_shapes = self._text_shapes(slide)
            family, variant, item_count, grid_shape, confidence = self._classify_slide(index, text_shapes)
            slots = self._build_slots(text_shapes)
            recipe = {
                "schema_version": "1.0",
                "slide_index": index,
                "name": f"{family}_{index}",
                "layout_family": family,
                "layout_variant": variant,
                "slots": slots,
                "bindings": {
                    slot["field"]: {
                        "shape_id": slot["shape_id"],
                        "bounds_emu": slot["bounds_emu"],
                    }
                    for slot in slots
                },
                "repeat_bindings": self._repeat_binding(item_count, grid_shape),
            }
            slide_recipes.append(recipe)
            catalog.append(
                {
                    "slide_index": index,
                    "layout_family": family,
                    "layout_variant": variant,
                    "item_count": item_count,
                    "grid_shape": grid_shape,
                    "confidence": confidence,
                }
            )
            colors, fonts, sizes = self._collect_style(slide)
            all_colors.extend(colors)
            all_fonts.extend(fonts)
            all_sizes.extend(sizes)

        design_system = self._design_system(prs, all_colors, all_fonts, all_sizes)
        grid_extendable = [
            item["slide_index"]
            for item in catalog
            if item["layout_family"] == "grid" and (item["item_count"] or 0) >= 3
        ]
        available = sorted({item["layout_family"] for item in catalog})
        schema = {
            "schema_version": "0.1",
            "generated_at": now_iso(),
            "source_pptx": str(template_path),
            "slide_dimensions_emu": [int(prs.slide_width), int(prs.slide_height)],
            "slides": slide_recipes,
            "design_system": design_system,
            "layout_catalog": catalog,
            "extension_capabilities": {
                "grid_extendable": grid_extendable,
                "sequence_extendable": [],
                "list_extendable": [],
            },
            "coverage_report": {
                "available_layout_families": available,
                "missing_common_layout_families": [item for item in COMMON_LAYOUTS if item not in available],
                "low_confidence_layouts": [
                    item["slide_index"] for item in catalog if item["confidence"] < 0.55
                ],
            },
        }
        return schema

    def _text_shapes(self, slide: Any) -> list[dict[str, Any]]:
        rows: list[dict[str, Any]] = []
        for shape in iter_shapes(slide.shapes):
            text = shape_text(shape)
            if not text:
                continue
            rows.append(
                {
                    "shape_id": int(shape.shape_id),
                    "text": text,
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                    "area": int(shape.width) * int(shape.height),
                }
            )
        rows.sort(key=lambda item: (item["top"], item["left"]))
        return rows

    def _classify_slide(
        self, index: int, text_shapes: list[dict[str, Any]]
    ) -> tuple[str, str, int | None, str | None, float]:
        if index == 1:
            return "cover", "title_cover", None, None, 0.9
        if len(text_shapes) <= 2:
            return "divider", "section_divider", None, None, 0.85

        numeric_labels = [
            item for item in text_shapes if re.fullmatch(r"\d{1,2}[.)]?", item["text"].strip())
        ]
        if len(numeric_labels) >= 3:
            labels = sorted(numeric_labels, key=lambda item: item["left"])
            xs = [item["left"] for item in labels]
            regular = self._is_regular(xs)
            count = len(labels)
            if regular:
                return "grid", f"{count}_item_grid", count, f"1x{count}", 0.8
            return "sequence", f"{count}_step_sequence", count, None, 0.62

        if any("|" in item["text"] for item in text_shapes) and len(text_shapes) >= 5:
            return "content", "structured_list", None, None, 0.6
        return "content", "text_content", None, None, 0.58

    def _is_regular(self, xs: list[int]) -> bool:
        if len(xs) < 3:
            return False
        gaps = [b - a for a, b in zip(xs, xs[1:])]
        avg = sum(gaps) / len(gaps)
        return all(abs(gap - avg) <= EMU_PER_INCH * 0.35 for gap in gaps)

    def _build_slots(self, text_shapes: list[dict[str, Any]]) -> list[dict[str, Any]]:
        if not text_shapes:
            return []
        title_candidate = max(text_shapes, key=lambda item: item["area"])
        slots: list[dict[str, Any]] = []
        for idx, shape in enumerate(text_shapes):
            if shape is title_candidate:
                field = "title"
            elif idx == 0:
                field = "eyebrow"
            elif idx == 1:
                field = "subtitle"
            elif idx == 2:
                field = "body"
            else:
                field = f"content_{idx - 1}"
            slots.append(
                {
                    "field": field,
                    "kind": "text",
                    "shape_id": shape["shape_id"],
                    "sample_text": shape["text"][:120],
                    "bounds_emu": [
                        shape["left"],
                        shape["top"],
                        shape["width"],
                        shape["height"],
                    ],
                }
            )
        return slots

    def _repeat_binding(self, item_count: int | None, grid_shape: str | None) -> dict[str, Any] | None:
        if not item_count:
            return None
        return {
            "slot_id": "grid_items",
            "kind": "grid",
            "source_count": item_count,
            "grid_shape": grid_shape,
            "item_fields": ["label", "title", "body"],
        }

    def _collect_style(self, slide: Any) -> tuple[list[str], list[str], list[int]]:
        colors: list[str] = []
        fonts: list[str] = []
        sizes: list[int] = []
        for shape in iter_shapes(slide.shapes):
            try:
                fill = shape.fill
                if fill and fill.fore_color and fill.fore_color.rgb:
                    colors.append(f"#{fill.fore_color.rgb}")
            except Exception:
                pass
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name:
                        fonts.append(run.font.name)
                    if run.font.size:
                        sizes.append(round(run.font.size.pt))
                    try:
                        if run.font.color and run.font.color.rgb:
                            colors.append(f"#{run.font.color.rgb}")
                    except Exception:
                        pass
        return colors, fonts, sizes

    def _design_system(
        self, prs: Presentation, colors: list[str], fonts: list[str], sizes: list[int]
    ) -> dict[str, Any]:
        palette = [color for color, _ in Counter(colors).most_common(8)]
        if not palette:
            palette = ["#174A7C", "#F2B705", "#FFFFFF", "#1F2933"]
        font = Counter(fonts).most_common(1)[0][0] if fonts else "Aptos"
        title_max = max(sizes) if sizes else 28
        body_min = min([size for size in sizes if size >= 8], default=14)
        width = int(prs.slide_width)
        height = int(prs.slide_height)
        return {
            "colors": {
                "primary": palette[0],
                "accent": palette[1] if len(palette) > 1 else palette[0],
                "background": "#FFFFFF",
                "text_primary": "#1F2933",
                "text_secondary": "#52606D",
                "palette": palette,
            },
            "typography": {
                "title": {"font": font, "size_range_pt": [18, max(24, title_max)], "weight": "bold"},
                "body": {"font": font, "size_range_pt": [12, max(14, body_min + 4)], "weight": "regular"},
                "caption": {"font": font, "size_range_pt": [10, 12], "weight": "regular"},
            },
            "spacing_emu": {
                "margin_left": int(width * 0.07),
                "margin_right": int(width * 0.07),
                "title_y": int(height * 0.08),
                "content_top": int(height * 0.29),
                "grid_gap": int(width * 0.018),
            },
            "shape_style": {
                "corner_radius": "small",
                "line_width_pt": 1.0,
                "shadow": "none",
            },
        }
